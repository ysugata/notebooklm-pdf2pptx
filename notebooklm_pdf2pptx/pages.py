"""入力の取得: PDF / 画像群 をページ毎の「キャンバス」へ正規化する。

キャンバス = そのページの基準ビットマップ(可能な限りネイティブ解像度)。
以後の全座標はキャンバスのピクセル空間で扱う。

PDFページが「全面1枚画像 + オーバーレイ(テキスト/小画像)」で構成される場合
(PowerPoint経由で書き出されたNotebookLMスライドが典型)は、
  - キャンバス: 埋め込み画像そのもの(再サンプリングなし)
  - オーバーレイテキスト: PDFのspan情報(フォント・サイズ・色が正確) → OCR不要・修復不要
  - オーバーレイ画像: 個別ピクチャとして再配置
に分解する。それ以外のページは高解像度レンダリングにフォールバックする。
"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path

import cv2
import fitz
import numpy as np

from .config import Settings

SUBSET_PREFIX = re.compile(r"^[A-Z]{6}\+")

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}


@dataclass
class OverlayText:
    text: str
    bbox_px: tuple[float, float, float, float]
    font: str
    size_pt_pdf: float        # PDF空間でのポイントサイズ(スライドptへは別途スケール)
    color_rgb: tuple[int, int, int]
    bold: bool
    italic: bool
    origin_y_px: float        # ベースラインy(キャンバスpx)


@dataclass
class OverlayImage:
    data: bytes               # エンコード済み画像 (PNG/JPEG)
    ext: str
    bbox_px: tuple[float, float, float, float]


@dataclass
class CanvasPage:
    number: int
    image: np.ndarray                     # BGR
    needs_ocr: bool
    px_per_pt: float | None = None        # PDF由来ページのみ (キャンバスpx / PDF pt)
    overlay_texts: list[OverlayText] = field(default_factory=list)
    overlay_images: list[OverlayImage] = field(default_factory=list)
    source_kind: str = "image"            # image | pdf-native | pdf-render | pptx
    slide_emu: tuple[int, int] | None = None   # PPTX入力時: 出力へ引き継ぐスライドサイズ
    native_xml: list[str] = field(default_factory=list)  # PPTX入力時: 持ち越すネイティブ図形XML
    cover_rects_px: list[tuple[float, float, float, float]] = field(default_factory=list)
    # ↑ 不透明な持ち越し要素が覆う領域。この下の焼き込み文字は原本で不可視のため
    #   復元・除去の対象から外す(隠しテキストの蘇生防止)

    @property
    def size(self) -> tuple[int, int]:
        h, w = self.image.shape[:2]
        return w, h


def _decode_image(data: bytes) -> np.ndarray | None:
    array = np.frombuffer(data, np.uint8)
    return cv2.imdecode(array, cv2.IMREAD_COLOR)


def _smask_is_opaque(doc: fitz.Document, smask_xref: int) -> bool:
    """実質不透明か (アンチエイリアス縁の数画素は許容する)。"""
    if not smask_xref:
        return True
    try:
        pix = fitz.Pixmap(doc, smask_xref)
        samples = np.frombuffer(pix.samples, np.uint8)
        return bool(np.mean(samples < 245) < 0.01)
    except Exception:
        return False


def _extract_overlay_image(doc: fitz.Document, xref: int, smask: int) -> tuple[bytes, str] | None:
    try:
        info = doc.extract_image(xref)
        if smask and not _smask_is_opaque(doc, smask):
            # ソフトマスクをアルファとして正確に合成する
            import io as _io

            from PIL import Image as _Image

            base = _Image.open(_io.BytesIO(info["image"])).convert("RGB")
            mask_info = doc.extract_image(smask)
            mask = _Image.open(_io.BytesIO(mask_info["image"])).convert("L")
            if mask.size != base.size:
                mask = mask.resize(base.size)
            base.putalpha(mask)
            buffer = _io.BytesIO()
            base.save(buffer, format="PNG")
            return buffer.getvalue(), "png"
        return info["image"], info["ext"]
    except Exception:
        return None


def _span_elements(page: fitz.Page, to_px) -> list[OverlayText]:
    elements: list[OverlayText] = []
    for block in page.get_text("dict").get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = str(span.get("text", ""))
                if not text.strip():
                    continue
                flags = int(span.get("flags", 0))
                x0, y0, x1, y1 = [float(v) for v in span["bbox"]]
                color = span.get("color", 0)
                rgb = ((color >> 16) & 255, (color >> 8) & 255, color & 255)
                elements.append(
                    OverlayText(
                        text=text,
                        bbox_px=(*to_px(x0, y0), *to_px(x1, y1)),
                        font=SUBSET_PREFIX.sub("", str(span.get("font", ""))),
                        size_pt_pdf=float(span.get("size", 12.0)),
                        color_rgb=rgb,
                        bold=bool(flags & 16) or "bold" in str(span.get("font", "")).lower(),
                        italic=bool(flags & 2),
                        origin_y_px=to_px(0, float(span.get("origin", (0, y1))[1]))[1],
                    )
                )
    return elements


def load_pdf_pages(pdf_path: Path, settings: Settings, pages: list[int] | None = None):
    """PDFをCanvasPageのジェネレータへ。100ページ級でもメモリ一定。"""
    doc = fitz.open(pdf_path)
    wanted = set(pages) if pages else None
    for page_index in range(doc.page_count):
        number = page_index + 1
        if wanted and number not in wanted:
            continue
        page = doc[page_index]
        yield _load_pdf_page(doc, page, number, settings)
    doc.close()


def _load_pdf_page(doc: fitz.Document, page: fitz.Page, number: int, settings: Settings) -> CanvasPage:
    page_rect = page.rect
    page_area = page_rect.width * page_rect.height
    # get_image_infoはsmaskを解決しない(常にNone)ため、get_imagesから対応表を作る
    smask_of = {img[0]: img[1] for img in page.get_images(full=True)}
    infos = page.get_image_info(xrefs=True)
    for info in infos:
        info["smask"] = smask_of.get(info.get("xref"), 0)
    base = None
    if infos:
        candidate = max(
            infos,
            key=lambda d: (d["bbox"][2] - d["bbox"][0]) * (d["bbox"][3] - d["bbox"][1]),
        )
        bbox = candidate["bbox"]
        coverage = ((bbox[2] - bbox[0]) * (bbox[3] - bbox[1])) / max(page_area, 1.0)
        if coverage >= settings.full_page_image_coverage:
            base = candidate

    if base is not None:
        canvas = _try_native_canvas(doc, page, base)
        if canvas is not None:
            image, place = canvas
            img_h, img_w = image.shape[:2]
            scale_x = img_w / max(place[2] - place[0], 1e-6)
            scale_y = img_h / max(place[3] - place[1], 1e-6)

            def to_px(x: float, y: float) -> tuple[float, float]:
                return ((x - place[0]) * scale_x, (y - place[1]) * scale_y)

            overlays_img: list[OverlayImage] = []
            for info in infos:
                if info["xref"] == base["xref"]:
                    continue
                extracted = _extract_overlay_image(doc, info["xref"], info.get("smask") or 0)
                if extracted is None:
                    continue
                data, ext = extracted
                b = info["bbox"]
                overlays_img.append(
                    OverlayImage(data=data, ext=ext, bbox_px=(*to_px(b[0], b[1]), *to_px(b[2], b[3])))
                )
            return CanvasPage(
                number=number,
                image=image,
                needs_ocr=True,
                px_per_pt=scale_x,
                overlay_texts=_span_elements(page, to_px),
                overlay_images=overlays_img,
                source_kind="pdf-native",
            )

    # フォールバック: ページ全体をレンダリング
    scale = settings.render_scale
    matrix = fitz.Matrix(scale, scale)
    pix = page.get_pixmap(matrix=matrix, alpha=False)
    image = np.frombuffer(pix.samples, np.uint8).reshape(pix.height, pix.width, pix.n)
    image = cv2.cvtColor(image, cv2.COLOR_RGB2BGR)

    def to_px_render(x: float, y: float) -> tuple[float, float]:
        return (x * scale, y * scale)

    has_visible_text = bool(page.get_text().strip())
    return CanvasPage(
        number=number,
        image=image,
        needs_ocr=not has_visible_text,
        px_per_pt=scale,
        overlay_texts=_span_elements(page, to_px_render) if has_visible_text else [],
        overlay_images=[],
        source_kind="pdf-render",
    )


def _try_native_canvas(doc: fitz.Document, page: fitz.Page, base: dict):
    """全面画像をネイティブ解像度で取得。透過が実質無い場合のみ。"""
    xref = base["xref"]
    smask = base.get("smask") or 0
    try:
        info = doc.extract_image(xref)
    except Exception:
        return None
    if smask and not _smask_is_opaque(doc, smask):
        return None
    image = _decode_image(info["image"])
    if image is None:
        return None
    # 変換行列に回転・せん断がある場合はフォールバック
    transform = base.get("transform")
    if transform and (abs(transform[1]) > 1e-3 or abs(transform[2]) > 1e-3):
        return None
    return image, tuple(float(v) for v in base["bbox"])


def load_image_pages(paths: list[Path], settings: Settings):
    for number, path in enumerate(paths, start=1):
        image = cv2.imread(str(path), cv2.IMREAD_COLOR)
        if image is None:
            raise FileNotFoundError(path)
        yield CanvasPage(number=number, image=image, needs_ocr=True, source_kind="image")


def resolve_input(input_path: Path) -> tuple[str, list[Path]]:
    """入力がPDFか画像(単体/ディレクトリ)かを判定する。"""
    if input_path.is_dir():
        images = sorted(
            p for p in input_path.iterdir() if p.suffix.lower() in IMAGE_SUFFIXES
        )
        if not images:
            raise FileNotFoundError(f"画像が見つかりません: {input_path}")
        return "images", images
    if input_path.suffix.lower() == ".pdf":
        return "pdf", [input_path]
    if input_path.suffix.lower() == ".pptx":
        return "pptx", [input_path]
    if input_path.suffix.lower() in IMAGE_SUFFIXES:
        return "images", [input_path]
    raise ValueError(f"未対応の入力形式: {input_path}")


# ---------------------------------------------------------------- PPTX入力
def _shape_abs_bbox_emu(shape, parents) -> tuple[float, float, float, float]:
    """グループ入れ子を遡って形状の絶対EMU座標を計算する(回転は非対応)。"""
    x = float(shape.left or 0)
    y = float(shape.top or 0)
    w = float(shape.width or 0)
    h = float(shape.height or 0)
    for group in reversed(parents):
        el = group._element
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        xfrm = el.find(".//a:xfrm", ns)
        ch_off = xfrm.find("a:chOff", ns)
        ch_ext = xfrm.find("a:chExt", ns)
        off = xfrm.find("a:off", ns)
        ext = xfrm.find("a:ext", ns)
        cox, coy = float(ch_off.get("x")), float(ch_off.get("y"))
        cex, cey = float(ch_ext.get("cx")), float(ch_ext.get("cy"))
        ox, oy = float(off.get("x")), float(off.get("y"))
        ex, ey = float(ext.get("cx")), float(ext.get("cy"))
        sx = ex / cex if cex else 1.0
        sy = ey / cey if cey else 1.0
        x = ox + (x - cox) * sx
        y = oy + (y - coy) * sy
        w *= sx
        h *= sy
    return x, y, w, h


def _iter_shapes_deep(shapes, parents=()):
    for shape in shapes:
        if shape.shape_type == 6:  # GROUP
            yield shape, parents
            yield from _iter_shapes_deep(shape.shapes, parents + (shape,))
        else:
            yield shape, parents


def _image_is_opaque(blob: bytes) -> bool:
    """オーバーレイ画像が実質不透明か(半透明スクリム等はカバー扱いしない)。"""
    array = np.frombuffer(blob, np.uint8)
    decoded = cv2.imdecode(array, cv2.IMREAD_UNCHANGED)
    if decoded is None:
        return False
    if decoded.ndim == 3 and decoded.shape[2] == 4:
        return bool(np.median(decoded[:, :, 3]) >= 230)
    return True


def _shape_has_opaque_fill(shape) -> bool:
    """ネイティブ図形が不透明塗りつぶしを持つか(XMLのspPr/solidFillで判定)。"""
    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
          "p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    element = shape._element
    sp_pr = element.find(".//p:spPr", ns) or element.find(".//a:spPr", ns)
    if sp_pr is None:
        sp_pr = element.find(
            "{http://schemas.openxmlformats.org/presentationml/2006/main}spPr")
    if sp_pr is None:
        # 名前空間を問わず spPr 直下を探す
        for child in element.iter():
            if child.tag.endswith("}spPr"):
                sp_pr = child
                break
    if sp_pr is None:
        return False
    fill = sp_pr.find("a:solidFill", ns)
    if fill is None:
        return False
    alpha = fill.find(".//a:alpha", ns)
    if alpha is not None and int(alpha.get("val", "100000")) < 90000:
        return False
    return True


def load_pptx_pages(pptx_path: Path, settings: Settings, pages: list[int] | None = None):
    """ほぼ全面画像で構成されたPPTX(NotebookLM出力を貼っただけのデッキ等)を
    CanvasPageへ分解する。

    - スライドを覆う最大のピクチャ → キャンバス(ネイティブ解像度のblob)
    - その他のピクチャ → OverlayImage(絶対座標で再配置)
    - ネイティブのテキスト/図形/グループ → XMLごと持ち越し
      (グループ内のピクチャはrelsが失われるためXMLから除去し、上記で個別再配置)
    - 出力スライドサイズは入力と同一にし、座標・サイズの変換を不要にする
    """
    import copy as _copy

    from lxml import etree
    from pptx import Presentation

    presentation = Presentation(str(pptx_path))
    slide_w = int(presentation.slide_width)
    slide_h = int(presentation.slide_height)
    slide_area = float(slide_w) * float(slide_h)
    wanted = set(pages) if pages else None

    PIC_TAG = "{http://schemas.openxmlformats.org/presentationml/2006/main}pic"

    for index, slide in enumerate(presentation.slides, start=1):
        if wanted and index not in wanted:
            continue
        base_shape = None
        base_bbox = None
        pics: list[tuple] = []
        for shape, parents in _iter_shapes_deep(slide.shapes):
            if shape.shape_type == 13:  # PICTURE
                bbox = _shape_abs_bbox_emu(shape, parents)
                coverage = (bbox[2] * bbox[3]) / slide_area
                pics.append((shape, bbox, coverage))
        if pics:
            best = max(pics, key=lambda t: t[2])
            if best[2] >= settings.full_page_image_coverage:
                base_shape, base_bbox = best[0], best[1]

        if base_shape is None:
            raise ValueError(
                f"slide {index}: 全面画像が見つかりません (PPTX入力は全面画像スライドのみ対応)")

        image = _decode_image(base_shape.image.blob)
        if image is None:
            raise ValueError(f"slide {index}: 画像をデコードできません")
        img_h, img_w = image.shape[:2]
        # キャンバスpx ↔ スライドEMU の変換 (基準はベース画像の配置範囲)
        scale_x = img_w / max(base_bbox[2], 1.0)
        scale_y = img_h / max(base_bbox[3], 1.0)

        def to_px_rect(bbox):
            x0 = (bbox[0] - base_bbox[0]) * scale_x
            y0 = (bbox[1] - base_bbox[1]) * scale_y
            return (x0, y0, x0 + bbox[2] * scale_x, y0 + bbox[3] * scale_y)

        overlays_img: list[OverlayImage] = []
        cover_rects: list[tuple[float, float, float, float]] = []
        for shape, bbox, _cov in pics:
            if shape is base_shape:
                continue
            rect = to_px_rect(bbox)
            overlays_img.append(OverlayImage(
                data=shape.image.blob, ext=shape.image.ext, bbox_px=rect))
            if _image_is_opaque(shape.image.blob):
                cover_rects.append(rect)

        # 塗りつぶし付きネイティブ図形もカバー領域として扱う
        for shape, parents in _iter_shapes_deep(slide.shapes):
            if shape.shape_type in (6, 13):  # GROUP/PICTUREは対象外(上で処理)
                continue
            if _shape_has_opaque_fill(shape):
                cover_rects.append(to_px_rect(_shape_abs_bbox_emu(shape, parents)))

        native_xml: list[str] = []
        for shape in slide.shapes:  # トップレベルのみ(グループは丸ごと)
            if shape.shape_type == 13:
                continue  # ピクチャは上で処理済み
            element = _copy.deepcopy(shape._element)
            # relsを持つ要素(ピクチャ・ハイパーリンク)を除去して破損を防ぐ
            for pic in element.findall(f".//{PIC_TAG}"):
                pic.getparent().remove(pic)
            for tag in ("hlinkClick", "hlinkHover", "hlinkMouseOver"):
                for node in element.findall(
                        ".//{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag):
                    node.getparent().remove(node)
            # ピクチャしか無かったグループ等、実体が残らない場合はスキップ
            has_content = element.findall(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}t") or                 shape.shape_type != 6
            if has_content:
                native_xml.append(etree.tostring(element, encoding="unicode"))

        px_per_pt = scale_x * 12700.0  # EMU→pt換算込み (px / pt)
        yield CanvasPage(
            number=index,
            image=image,
            needs_ocr=True,
            px_per_pt=px_per_pt,
            overlay_texts=[],
            overlay_images=overlays_img,
            source_kind="pptx",
            slide_emu=(slide_w, slide_h),
            native_xml=native_xml,
            cover_rects_px=cover_rects,
        )
