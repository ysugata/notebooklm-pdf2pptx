"""PPTX組み立て: 背景画像 + 編集可能テキスト。

位置決めの原理:
  ソルバーが解いた「ベースラインのキャンバス座標」を基準に、PowerPointの
  行レイアウト(フォントメトリクス由来の行高とベースライン位置)を逆算して
  テキストボックスの座標を決める。整列(左/中央/右)に応じてボックスを
  インク実測に合わせてアンカーするため、編集時も自然に振る舞う。
"""
from __future__ import annotations

import json
from pathlib import Path

import cv2
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from .config import Settings, JP_FALLBACK_CHAIN
from .fontlib import FontLibrary, line_metrics_units

EMU_PER_IN = 914400
SLIDE_W_EMU = 12192000  # 13.333in
SLIDE_H_EMU = 6858000   # 7.5in


# ---------------------------------------------------------------- font metrics
def _win_metrics_pt(font_path: str, font_index: int, size_pt: float) -> tuple[float, float]:
    """PowerPointの行計算に使う (ascent_pt, descent_pt)。

    Windows系レンダラはOS/2 usWinAscent/usWinDescentを使う。
    fsSelection USE_TYPO_METRICS が立つフォントは sTypo* を使う。
    """
    units = line_metrics_units(font_path, font_index)
    upem = units["upem"]
    if units.get("use_typo") and "typo_ascent" in units:
        ascent = units["typo_ascent"] + units.get("typo_linegap", 0) / 2
        descent = units["typo_descent"] + units.get("typo_linegap", 0) / 2
    elif "win_ascent" in units:
        ascent = units["win_ascent"]
        descent = units["win_descent"]
    else:
        ascent = units.get("hhea_ascent", upem * 0.8)
        descent = units.get("hhea_descent", upem * 0.2)
    return ascent / upem * size_pt, descent / upem * size_pt


def first_baseline_offset_pt(font_path: str, font_index: int, size_pt: float,
                             exact_spacing_pt: float | None) -> float:
    """テキストボックス上端(マージン0)から最初のベースラインまでの距離。

    行高固定(spcPts)時のLibreOffice実測則(tools/calibrate2.py、16ケースの
    PDFベースライン座標で±0.03pt一致): baseline = 行高 − 0.2×フォントサイズ。
    ディセント扱いはフォントメトリクス非依存の固定20%。
    """
    if exact_spacing_pt is None:
        ascent, _descent = _win_metrics_pt(font_path, font_index, size_pt)
        return ascent
    return exact_spacing_pt - 0.2 * size_pt


def natural_line_height_pt(font_path: str, font_index: int, size_pt: float) -> float:
    ascent, descent = _win_metrics_pt(font_path, font_index, size_pt)
    return ascent + descent


# ---------------------------------------------------------------- XML helpers
# CT_TextCharacterProperties の子要素スキーマ順 (ECMA-376)
_RPR_ORDER = (
    "a:ln", "a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill",
    "a:grpFill", "a:effectLst", "a:effectDag", "a:highlight", "a:uLnTx", "a:uLn",
    "a:uFillTx", "a:uFill", "a:latin", "a:ea", "a:cs", "a:sym",
    "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst",
)


def _insert_ordered(rpr, element) -> None:
    """rPrへスキーマ順を守って子要素を挿入する (順序違反はKeynote等が読めない)。"""
    tag_index = {qn(tag): i for i, tag in enumerate(_RPR_ORDER)}
    own = tag_index.get(element.tag, len(_RPR_ORDER))
    for child in rpr:
        if tag_index.get(child.tag, -1) > own:
            child.addprevious(element)
            return
    rpr.append(element)


def _get_or_add_ordered(rpr, tag: str):
    node = rpr.find(qn(tag))
    if node is None:
        node = OxmlElement(tag)
        _insert_ordered(rpr, node)
    return node


def _set_run_fonts(run, typeface: str) -> None:
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        _get_or_add_ordered(rpr, tag).set("typeface", typeface)


def _set_char_spacing(run, spacing_pt: float) -> None:
    if abs(spacing_pt) < 0.05:
        return
    run._r.get_or_add_rPr().set("spc", str(int(round(spacing_pt * 100))))


def _disable_kerning(run) -> None:
    # kern属性は「ペアカーニングを有効にする最小フォントサイズ」(1/100pt)。
    # "0"の意味は実装依存で曖昧なため、実用上あり得ない大きさ(100pt)を
    # 指定してカーニングを実質無効化する。照合レンダリング(PILは非カーニング)
    # と送り幅を一致させるため。
    run._r.get_or_add_rPr().set("kern", "10000")


def _set_gradient_fill(run, top_rgb, bottom_rgb) -> None:
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
        node = rpr.find(qn(tag))
        if node is not None:
            rpr.remove(node)
    grad = OxmlElement("a:gradFill")
    gs_list = OxmlElement("a:gsLst")
    for pos, rgb in (("0", top_rgb), ("100000", bottom_rgb)):
        gs = OxmlElement("a:gs")
        gs.set("pos", pos)
        color = OxmlElement("a:srgbClr")
        color.set("val", "".join(f"{int(v):02X}" for v in rgb))
        gs.append(color)
        gs_list.append(gs)
    grad.append(gs_list)
    lin = OxmlElement("a:lin")
    lin.set("ang", "5400000")  # 90° = 上→下
    lin.set("scaled", "0")
    grad.append(lin)
    _insert_ordered(rpr, grad)


def _add_glow(run, radius_pt: float, rgb, alpha_pct: int = 55) -> None:
    rpr = run._r.get_or_add_rPr()
    effects = _get_or_add_ordered(rpr, "a:effectLst")
    glow = OxmlElement("a:glow")
    glow.set("rad", str(int(round(radius_pt * 12700))))
    color = OxmlElement("a:srgbClr")
    color.set("val", "".join(f"{int(v):02X}" for v in rgb))
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(alpha_pct * 1000))
    color.append(alpha)
    glow.append(color)
    effects.append(glow)


def _set_exact_line_spacing(paragraph, spacing_pt: float) -> None:
    ppr = paragraph._p.get_or_add_pPr()
    ln_spc = ppr.find(qn("a:lnSpc"))
    if ln_spc is None:
        ln_spc = OxmlElement("a:lnSpc")
        ppr.insert(0, ln_spc)
    else:
        for child in list(ln_spc):
            ln_spc.remove(child)
    spc_pts = OxmlElement("a:spcPts")
    spc_pts.set("val", str(int(round(spacing_pt * 100))))
    ln_spc.append(spc_pts)


# ---------------------------------------------------------------- main builder
def build_presentation(processed: list[dict], pages_dir: Path, output_path: Path,
                       settings: Settings, library: FontLibrary) -> dict:
    presentation = Presentation()
    first = min(processed, key=lambda r: r["number"])
    if first.get("slide_emu"):
        # PPTX入力: 出力サイズを入力スライドに合わせる
        # (持ち越すネイティブ図形の座標・サイズをそのまま使えるようにする)
        slide_w_emu, slide_h_emu = first["slide_emu"]
    else:
        slide_w_emu = SLIDE_W_EMU
        # スライド高さ: 16:9近傍(±2%)なら16:9へ正規化、それ以外は元アスペクト維持
        cw, ch = first["canvas_size"]
        slide_h_emu = SLIDE_H_EMU
        if settings.force_16_9:
            aspect = cw / ch
            if abs(aspect - 16 / 9) / (16 / 9) > 0.02:
                slide_h_emu = int(round(SLIDE_W_EMU * ch / cw))
        else:
            slide_h_emu = int(round(SLIDE_W_EMU * ch / cw))
    presentation.slide_width = Emu(slide_w_emu)
    presentation.slide_height = Emu(slide_h_emu)
    blank = presentation.slide_layouts[6]

    n_text = 0
    for record in sorted(processed, key=lambda r: r["number"]):
        page_dir = pages_dir / f"{record['number']:03d}"
        layout = json.loads((page_dir / "layout.json").read_text(encoding="utf-8"))
        slide = presentation.slides.add_slide(blank)
        canvas_w, canvas_h = layout["canvas_size"]
        emu_per_px_x = slide_w_emu / canvas_w
        emu_per_px_y = slide_h_emu / canvas_h
        pt_per_px = layout["pt_per_px"]

        background_path = _background_for_slide(page_dir, settings)
        picture = slide.shapes.add_picture(
            str(background_path), 0, 0, presentation.slide_width, presentation.slide_height
        )
        picture.name = f"Background page {record['number']}"

        # Z順序: 背景 → 復元テキスト(焼き込みの置換なので背景直上) →
        # オーバーレイ画像 → オーバーレイテキスト → ネイティブ持ち越し図形
        for block in layout["blocks"]:
            _add_block(slide, block, emu_per_px_x, emu_per_px_y, pt_per_px)
            n_text += len(block["lines"])

        for overlay in record.get("overlay_images", []):
            x0, y0, x1, y1 = overlay["bbox"]
            slide.shapes.add_picture(
                str(page_dir / overlay["file"]),
                Emu(int(x0 * emu_per_px_x)), Emu(int(y0 * emu_per_px_y)),
                Emu(int((x1 - x0) * emu_per_px_x)), Emu(int((y1 - y0) * emu_per_px_y)),
            )

        for overlay in layout.get("overlay_texts", []):
            _add_overlay_text(slide, overlay, emu_per_px_x, emu_per_px_y, pt_per_px,
                              library)
            n_text += 1

        for xml_name in record.get("native_xml", []):
            from lxml import etree
            xml_path = page_dir / xml_name
            if xml_path.is_file():
                element = etree.fromstring(xml_path.read_text(encoding="utf-8"))
                slide.shapes._spTree.append(element)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return {
        "output": str(output_path),
        "pages": len(processed),
        "text_elements": n_text,
        "reviews": {str(r["number"]): r["n_review"] for r in processed if r.get("n_review")},
    }


def _background_for_slide(page_dir: Path, settings: Settings) -> Path:
    """背景をJPEG化してPPTXサイズを抑える(100ページ対応)。"""
    png = page_dir / "background.png"
    jpg = page_dir / "background.jpg"
    if not jpg.is_file() or jpg.stat().st_mtime < png.stat().st_mtime:
        image = cv2.imread(str(png))
        cv2.imwrite(str(jpg), image, [cv2.IMWRITE_JPEG_QUALITY, settings.jpeg_quality])
    return jpg


def _add_block(slide, block: dict, emu_x: float, emu_y: float, pt_per_px: float) -> None:
    lines = block["lines"]
    first = lines[0]
    align = block["align"]
    n = len(lines)

    size_pt = first["size_pt"]
    font_path, font_index = first["font_path"], first["font_index"]

    # 常に行高を固定(spcPts)する: 通常行送りはレンダラ間でメトリクス解釈が
    # 異なる(CoreText=typo/hhea, GDI=win)が、行高固定ならベースラインは
    # 「行高 - ディセント」で安定する (tools/calibrate.py で±3px以内を確認)。
    line_h_pt = block["pitch_px"] * pt_per_px if n > 1 else natural_line_height_pt(
        font_path, font_index, size_pt)
    pitch_pt = line_h_pt
    fbo_pt = first_baseline_offset_pt(font_path, font_index, size_pt, line_h_pt)

    baseline0_px = first["baseline_y_px"]
    top_px = baseline0_px - fbo_pt / pt_per_px
    # 行間は一律ピッチではなく実測ベースライン間隔を段落ごとに指定する。
    # 中央値ピッチの一律適用は±1px程度の実測揺らぎを累積させ、
    # 行数が多いブロックで下の行ほど位置がずれるため。
    baselines = [l["baseline_y_px"] for l in lines]
    gaps_pt = [(baselines[i] - baselines[i - 1]) * pt_per_px for i in range(1, n)]
    height_pt = fbo_pt + sum(gaps_pt) + max(
        _win_metrics_pt(font_path, font_index, size_pt)[1], 2.0)

    # 幅: 最長行の送り幅 + 余裕。整列に応じて左端を決める。
    widths_px = [l["advance_w_px"] for l in lines]
    max_w_px = max(widths_px)
    # wrap=noneだが、wrap=none非対応レンダラ(QuickLook等)が折り返さないよう
    # 幅に十分な余裕を持たせる。整列アンカーは実測インク基準なので
    # 幅を広げても文字位置はずれない。
    box_w_px = max_w_px * 1.10 + 16
    if align == "center":
        centers = [(l["origin_x_px"] + l["advance_w_px"] / 2) for l in lines]
        center_px = sum(centers) / len(centers)
        left_px = center_px - box_w_px / 2
    elif align == "right":
        rights = [l["origin_x_px"] + l["advance_w_px"] for l in lines]
        right_px = max(rights)
        left_px = right_px - box_w_px
    else:
        left_px = min(l["origin_x_px"] for l in lines)

    textbox = slide.shapes.add_textbox(
        Emu(int(left_px * emu_x)),
        Emu(int(top_px * emu_y)),
        Emu(int(box_w_px * emu_x)),
        Emu(int(height_pt / pt_per_px * emu_y)),
    )
    textbox.name = f"Text: {first['text'][:28]}"
    frame = textbox.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.word_wrap = False
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
        }[align]
        paragraph.space_before = paragraph.space_after = Pt(0)
        _set_exact_line_spacing(
            paragraph, pitch_pt if index == 0 else gaps_pt[index - 1])
        glow = line.get("glow")

        def emit_run(text: str, color, gradient) -> None:
            run = paragraph.add_run()
            run.text = text
            run.font.size = Pt(line["size_pt"])
            run.font.bold = bool(line["bold"])
            _set_run_fonts(run, line["font_family"])
            _set_char_spacing(run, line["char_spacing_pt"])
            _disable_kerning(run)
            if gradient:
                run.font.color.rgb = RGBColor(*gradient[0])
                _set_gradient_fill(run, gradient[0], gradient[1])
            else:
                run.font.color.rgb = RGBColor(*color)
            if glow:
                _add_glow(run, glow["radius_pt"], glow["color"])

        color_runs = line.get("color_runs")
        if color_runs and len(color_runs) > 1:
            # 行内の部分強調: 書式は共有し色だけをランごとに変える。
            # 横方向の色変化はグラデーション(縦)とは別概念なので適用しない。
            for start, end, rgb in color_runs:
                segment = line["text"][start:end]
                if segment:
                    emit_run(segment, rgb, None)
        else:
            emit_run(line["text"], line["color"], line.get("gradient"))


def _add_overlay_text(slide, overlay: dict, emu_x: float, emu_y: float,
                      pt_per_px: float, library: FontLibrary) -> None:
    """PDF由来オーバーレイ: フォント・サイズ・色が正確に分かっている要素。"""
    from .fontlib import render_metrics, REF_SIZE

    x0, y0, x1, y1 = overlay["bbox_px"]
    text = overlay["text"]
    # PDF pt → キャンバスpx → スライドpt
    size_pt = overlay["size_pt_pdf"] * overlay.get("px_per_pt", 1.0) * pt_per_px

    family = _map_pdf_font(overlay["font"], library, text)
    face = library.face(family, bold=overlay.get("bold", False))
    baseline_px = overlay.get("origin_y_px", y1)

    spacing_pt = 0.0
    if face is not None:
        rm = render_metrics(text, face.path, face.index)
        if rm is not None and rm.n_gaps:
            scale = (size_pt / pt_per_px) / REF_SIZE
            natural_w = rm.adv_w * scale
            spacing_pt = ((x1 - x0) - natural_w) / rm.n_gaps * pt_per_px
            spacing_pt = max(-12.0, min(12.0, spacing_pt))

    if face is not None:
        line_h_pt = natural_line_height_pt(face.path, face.index, size_pt)
        fbo_pt = first_baseline_offset_pt(face.path, face.index, size_pt, line_h_pt)
    else:
        line_h_pt = size_pt * 1.2
        fbo_pt = size_pt
    top_px = baseline_px - fbo_pt / pt_per_px
    height_pt = fbo_pt + size_pt * 0.35

    textbox = slide.shapes.add_textbox(
        Emu(int(x0 * emu_x)), Emu(int(top_px * emu_y)),
        Emu(int((x1 - x0) * 1.05 * emu_x)), Emu(int(height_pt / pt_per_px * emu_y)),
    )
    textbox.name = f"Overlay text: {text[:28]}"
    frame = textbox.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.word_wrap = False
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    _set_exact_line_spacing(paragraph, line_h_pt)
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(round(size_pt, 2))
    run.font.bold = bool(overlay.get("bold"))
    run.font.color.rgb = RGBColor(*overlay["color"])
    _set_run_fonts(run, family)
    _set_char_spacing(run, spacing_pt)
    _disable_kerning(run)


def _map_pdf_font(pdf_font: str, library: FontLibrary, text: str) -> str:
    """PDFのフォント名をインストール済みfamilyへマップする。"""
    name = pdf_font.replace("-", " ").replace("_", " ")
    for suffix in (" Bold", " Italic", " Regular", " Medium"):
        if name.endswith(suffix):
            name = name[: -len(suffix)]
    candidates = [pdf_font, name, name.replace(" ", "")]
    # 例: "BIZ-UDPGothic" → "BIZ UDPGothic"
    for candidate in candidates:
        if library.find(candidate):
            return candidate
    for fallback in JP_FALLBACK_CHAIN:
        if library.find(fallback):
            return fallback
    return name
