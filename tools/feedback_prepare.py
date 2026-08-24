#!/usr/bin/env python3
"""修正指示(注記)つきPPTXから、修正タスク一覧 tasks.json を生成する。

対象は編集可能化済みのPPTX(注記図形が持ち越されているもの)。
やること(すべて決定論):
  1. rules/annotations.yaml に従って注記図形を検出・分類 (A=文字化け /
     B=置換指示 / C=要人間判断)
  2. 各注記の指し先(最寄りのテキスト図形)を特定し、現在の本文を記録
  3. A分類には辞書検証つき修復(textfix)の提案を suggested として同梱
  4. LibreOfficeで該当スライドをレンダリングし、注記周辺の切り抜き画像を出力
  5. tasks.json(+ crops/)を書き出す。ハッシュを埋め込み、apply側が検証する

エージェント(Codex/Claude/人間)の仕事は answers.json を書くことだけ。
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

import yaml  # noqa: E402
from pptx import Presentation  # noqa: E402
from notebooklm_pdf2pptx.pages import _iter_shapes_deep, _shape_abs_bbox_emu  # noqa: E402


def load_rules() -> dict:
    return yaml.safe_load((ROOT / "rules" / "annotations.yaml").read_text("utf-8"))


def shape_geometry(shape) -> str:
    el = shape._element
    geom = el.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
    return geom.get("prst") if geom is not None else ""


def is_marker(shape, rules: dict) -> bool:
    m = rules["marker"]
    text = shape.text_frame.text.strip() if getattr(shape, "has_text_frame", False) else ""
    if len(text) > m["max_text_chars"]:
        return False
    geom = shape_geometry(shape)
    if any(geom.startswith(p) for p in m["geometry_prefixes"]):
        return True
    return bool(text) and any(k in text for k in m["instruction_keywords"])


def classify(text: str, rules: dict) -> str:
    c = rules["classify"]
    if any(k in text for k in c["garble_keywords"]):
        return "A"
    if any(p in text for p in c["replacement_patterns"]):
        return "B"
    return "C"


def rect_of(shape):
    return (shape.left or 0, shape.top or 0,
            (shape.left or 0) + (shape.width or 0),
            (shape.top or 0) + (shape.height or 0))


def rect_distance(a, b) -> float:
    dx = max(b[0] - a[2], a[0] - b[2], 0)
    dy = max(b[1] - a[3], a[1] - b[3], 0)
    return (dx * dx + dy * dy) ** 0.5


def callout_tip(shape, rect):
    """吹き出しの尻尾の先端座標(EMU、絶対)。wedge系はadj1/adj2が中心からの
    比率で先端位置を持つ。取れなければNone。rectは図形の絶対座標。"""
    A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    gds = {g.get("name"): g.get("fmla")
           for g in shape._element.findall(f".//{A}avLst/{A}gd")}

    def val(name):
        f = gds.get(name)
        if f and f.startswith("val "):
            try:
                return float(f.split()[1]) / 100000.0
            except ValueError:
                return None
        return None

    a1, a2 = val("adj1"), val("adj2")
    if a1 is None or a2 is None:
        return None
    w = rect[2] - rect[0]
    h = rect[3] - rect[1]
    return (rect[0] + w * (0.5 + a1), rect[1] + h * (0.5 + a2))


def point_rect_distance(pt, rect) -> float:
    dx = max(rect[0] - pt[0], pt[0] - rect[2], 0)
    dy = max(rect[1] - pt[1], pt[1] - rect[3], 0)
    return (dx * dx + dy * dy) ** 0.5


def paragraphs_text(shape) -> list[str]:
    return ["".join(r.text for r in p.runs) for p in shape.text_frame.paragraphs]


def render_slides(pptx: Path, out_dir: Path) -> dict[int, Path]:
    """LibreOfficeで全スライドをPNG化する。失敗したら空dict(クロップなし運用)。"""
    out_dir.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory() as td:
        try:
            subprocess.run(["soffice", "--headless", "--convert-to", "pdf",
                            "--outdir", td, str(pptx)],
                           check=True, capture_output=True, timeout=600)
            pdf = next(Path(td).glob("*.pdf"))
            import fitz
            doc = fitz.open(pdf)
            pages = {}
            for i, page in enumerate(doc, 1):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                p = out_dir / f"slide_{i:03d}.png"
                pix.save(p)
                pages[i] = p
            return pages
        except Exception as exc:  # noqa: BLE001
            print(f"レンダリング不可(クロップなしで続行): {exc}")
            return {}


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--out-dir", type=Path, default=ROOT / "feedback_work")
    ap.add_argument("--work-dir", type=Path, default=None,
                    help="変換時のworkディレクトリ。指定すると除去前の原本画像"
                         "(pages/NNN/source.png)からも切り抜きを作る(判断精度向上)")
    args = ap.parse_args()

    rules = load_rules()
    prs = Presentation(str(args.pptx))
    out = args.out_dir
    crops_dir = out / "crops"
    crops_dir.mkdir(parents=True, exist_ok=True)

    try:
        from notebooklm_pdf2pptx.textfix import TextRepairer
        repairer = TextRepairer()
    except Exception:
        repairer = None

    tasks = []
    slides_with_tasks: set[int] = set()
    for idx, slide in enumerate(prs.slides, 1):
        # グループ入れ子も含めて全図形を走査し、絶対座標を持たせる
        deep = [(s, _shape_abs_bbox_emu(s, parents))
                for s, parents in _iter_shapes_deep(slide.shapes)
                if s.shape_type != 6]
        abs_rect = {id(s): (b[0], b[1], b[0] + b[2], b[1] + b[3]) for s, b in deep}
        shapes = [s for s, _b in deep]
        markers = [s for s in shapes if is_marker(s, rules)]
        texty = [s for s in shapes
                 if getattr(s, "has_text_frame", False)
                 and s.text_frame.text.strip() and s not in markers]
        for mk in markers:
            mtext = mk.text_frame.text.strip() if getattr(mk, "has_text_frame", False) else ""
            klass = classify(mtext, rules)
            # 指し先: 吹き出しの尻尾の先端が指すテキスト図形(先端座標が
            # 取れない図形は、注記に最も近いテキスト図形へフォールバック)
            tip = callout_tip(mk, abs_rect[id(mk)])
            if tip is not None and texty:
                target = min(texty,
                             key=lambda s: point_rect_distance(tip, abs_rect[id(s)]))
            else:
                target = min(texty,
                             key=lambda s: rect_distance(abs_rect[id(mk)],
                                                         abs_rect[id(s)]),
                             default=None)
            entry = {
                "id": f"s{idx:03d}_m{mk.shape_id}",
                "slide": idx,
                "class": klass,
                "marker_text": mtext,
                "marker_shape_id": mk.shape_id,
                "target": None,
                "suggested": None,
                "status": "open" if klass in ("A", "B") else "needs_human",
            }
            rect = list(abs_rect[id(mk)])
            if target is not None:
                paras = paragraphs_text(target)
                entry["target"] = {
                    "shape_id": target.shape_id,
                    "name": target.name,
                    "paragraphs": paras,
                }
                tr = abs_rect[id(target)]
                rect = [min(rect[0], tr[0]), min(rect[1], tr[1]),
                        max(rect[2], tr[2]), max(rect[3], tr[3])]
                if klass == "A" and repairer is not None and repairer.available:
                    fixed = [repairer.apply(p)[0] for p in paras]
                    if fixed != paras:
                        entry["suggested"] = fixed
            entry["_rect"] = rect
            tasks.append(entry)
            slides_with_tasks.add(idx)

    # クロップ生成
    if slides_with_tasks:
        pages = render_slides(args.pptx, out / "renders")
        sw = prs.slide_width
        sh = prs.slide_height
        if pages:
            import cv2
            for t in tasks:
                img_path = pages.get(t["slide"])
                if img_path is None:
                    continue
                img = cv2.imread(str(img_path))
                ih, iw = img.shape[:2]
                x0, y0, x1, y1 = t["_rect"]
                pad = int(0.03 * sw)
                px0 = max(0, int((x0 - pad) / sw * iw))
                py0 = max(0, int((y0 - pad) / sh * ih))
                px1 = min(iw, int((x1 + pad) / sw * iw))
                py1 = min(ih, int((y1 + pad) / sh * ih))
                if px1 - px0 > 8 and py1 - py0 > 8:
                    cp = crops_dir / f"{t['id']}.png"
                    cv2.imwrite(str(cp), img[py0:py1, px0:px1])
                    t["crop"] = str(cp.relative_to(out))
                # 除去前の原本画像からの切り抜き(最重要の判断根拠)。
                # 編集可能化後の描画はOCRの推測を写しただけだが、
                # 原本ピクセルには崩れ字の実際の形が残っている。
                if args.work_dir is not None:
                    src_png = args.work_dir / "pages" / f"{t['slide']:03d}" / "source.png"
                    if src_png.is_file():
                        simg = cv2.imread(str(src_png))
                        if simg is not None:
                            sh_, sw_ = simg.shape[:2]
                            sx0 = max(0, int((x0 - pad) / sw * sw_))
                            sy0 = max(0, int((y0 - pad) / sh * sh_))
                            sx1 = min(sw_, int((x1 + pad) / sw * sw_))
                            sy1 = min(sh_, int((y1 + pad) / sh * sh_))
                            if sx1 - sx0 > 8 and sy1 - sy0 > 8:
                                scp = crops_dir / f"{t['id']}_source.png"
                                # 小さい字も読めるよう2倍に拡大して保存
                                crop2 = cv2.resize(simg[sy0:sy1, sx0:sx1],
                                                   (0, 0), fx=2, fy=2,
                                                   interpolation=cv2.INTER_LANCZOS4)
                                cv2.imwrite(str(scp), crop2)
                                t["source_crop"] = str(scp.relative_to(out))

    for t in tasks:
        t.pop("_rect", None)
    payload = {
        "source": str(args.pptx),
        "source_sha256": hashlib.sha256(args.pptx.read_bytes()).hexdigest(),
        "tasks": tasks,
    }
    body = json.dumps(payload, ensure_ascii=False, indent=1)
    tasks_path = out / "tasks.json"
    tasks_path.write_text(body, "utf-8")
    digest = hashlib.sha256(body.encode()).hexdigest()
    (out / "tasks.sha256").write_text(digest, "utf-8")
    by_class = {}
    for t in tasks:
        by_class[t["class"]] = by_class.get(t["class"], 0) + 1
    print(f"tasks.json: {len(tasks)}件 (分類: {by_class}) -> {tasks_path}")
    print(f"hash: {digest[:16]}…  (apply側が検証します)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
