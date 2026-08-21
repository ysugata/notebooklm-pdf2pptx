"""ベースライン較正: 行レイアウト計算式を実レンダラで検証する。

既知のターゲット位置に文字を置いたPPTXを生成し、レンダリング結果の
インク位置を実測して、first_baseline_offset_pt の系統誤差を報告する。

使い方:
  .venv/bin/python tools/calibrate.py --out-dir work/calib
"""
from __future__ import annotations

import argparse
import subprocess
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import cv2
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from notebooklm_pdf2pptx.fontlib import FontLibrary, render_metrics, REF_SIZE
from notebooklm_pdf2pptx.pptx_writer import (
    SLIDE_H_EMU, SLIDE_W_EMU, _set_char_spacing, _set_exact_line_spacing,
    _set_run_fonts, _disable_kerning, first_baseline_offset_pt,
)
from notebooklm_pdf2pptx.config import Settings

CANVAS_W, CANVAS_H = 1376, 768
PT_PER_PX = 960.0 / CANVAS_W


def build_cases(library: FontLibrary) -> list[dict]:
    cases = []
    specs = [
        # spcPts固定行高 (行高 = size * 1.35)
        ("Noto Sans JP", True, "日本語ベースライン検証テキスト", 32.0, 0.0, 43.2),
        ("Noto Sans JP", False, "小さめの日本語テキスト検証", 14.0, 0.0, 18.9),
        ("Oswald", False, "LATIN BASELINE CHECK", 32.0, 0.0, 43.2),
        ("Oswald", True, "CONDENSED BOLD 24PT", 24.0, 2.0, 32.4),
        ("Hiragino Sans W6", False, "ヒラギノ検証テキスト", 28.0, -2.0, 37.8),
        ("Noto Sans JP", True, "行間固定の複数行検証\n二行目のテキスト", 20.0, 0.0, 34.0),
        ("Futura Condensed Medium", False, "FUTURA CONDENSED CHECK", 28.0, 0.0, 37.8),
    ]
    y_px = 60.0
    for family, bold, text, size_pt, spc, pitch_pt in specs:
        face = library.face(family, bold=bold)
        if face is None:
            continue
        cases.append({
            "face": face, "text": text, "size_pt": size_pt, "spc": spc,
            "pitch_pt": pitch_pt, "target_baseline_px": y_px, "left_px": 80.0,
        })
        n_lines = text.count("\n") + 1
        y_px += (size_pt / PT_PER_PX) * 1.6 * n_lines + 30
    return cases


def build_pptx(cases: list[dict], path: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = Emu(SLIDE_W_EMU)
    presentation.slide_height = Emu(SLIDE_H_EMU)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    emu_x = SLIDE_W_EMU / CANVAS_W
    emu_y = SLIDE_H_EMU / CANVAS_H

    for case in cases:
        face = case["face"]
        size_pt = case["size_pt"]
        fbo = first_baseline_offset_pt(face.path, face.index, size_pt, case["pitch_pt"])
        top_px = case["target_baseline_px"] - fbo / PT_PER_PX
        box = slide.shapes.add_textbox(
            Emu(int(case["left_px"] * emu_x)), Emu(int(top_px * emu_y)),
            Emu(int(1100 * emu_x)), Emu(int((size_pt * 3) / PT_PER_PX * emu_y)),
        )
        frame = box.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
        frame.word_wrap = False
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.vertical_anchor = MSO_ANCHOR.TOP
        for i, line in enumerate(case["text"].split("\n")):
            paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_before = paragraph.space_after = Pt(0)
            if case["pitch_pt"]:
                _set_exact_line_spacing(paragraph, case["pitch_pt"])
            run = paragraph.add_run()
            run.text = line
            run.font.size = Pt(size_pt)
            run.font.bold = face.bind_bold
            run.font.color.rgb = RGBColor(0, 0, 0)
            _set_run_fonts(run, face.typeface)
            _set_char_spacing(run, case["spc"])
            _disable_kerning(run)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(path))


def measure(cases: list[dict], render_path: Path) -> None:
    image = cv2.imread(str(render_path), cv2.IMREAD_GRAYSCALE)
    image = cv2.resize(image, (CANVAS_W, CANVAS_H))
    binary = (image < 128).astype(np.uint8)
    print(f"{'case':44s} {'期待bl':>7s} {'実測bl':>7s} {'差px':>6s}")
    for case in cases:
        face = case["face"]
        text = case["text"].split("\n")[0]
        rm = render_metrics(text, face.path, face.index)
        if rm is None:
            continue
        scale = (case["size_pt"] / PT_PER_PX) / REF_SIZE
        # 期待インク上端 = baseline - (ascent - ink_top_from_ascent相当)
        expected_top = case["target_baseline_px"] - (rm.ascent - rm.ink_top_from_ascent) * scale
        expected_bottom = expected_top + rm.ink_h * scale
        y0 = max(0, int(expected_top) - 14)
        y1 = min(CANVAS_H, int(expected_bottom) + 14)
        strip = binary[y0:y1, int(case["left_px"]) - 5: int(case["left_px"]) + 700]
        ys, _xs = np.nonzero(strip)
        if len(ys) == 0:
            print(f"{text[:40]:44s} ink not found")
            continue
        measured_top = y0 + ys.min()
        measured_baseline = measured_top + (rm.ascent - rm.ink_top_from_ascent) * scale - (
            rm.ascent - rm.ink_top_from_ascent) * scale + (
            case["target_baseline_px"] - expected_top)
        # 単純化: インク上端の期待/実測差 = ベースライン差
        delta = measured_top - expected_top
        label = f"{face.typeface} {case['size_pt']}pt" + (
            f" pitch={case['pitch_pt']}" if case["pitch_pt"] else "")
        print(f"{label:44s} {case['target_baseline_px']:7.1f} {case['target_baseline_px'] + delta:7.1f} {delta:+6.1f}")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out-dir", type=Path, default=Path("work/calib"))
    args = parser.parse_args()
    settings = Settings()
    library = FontLibrary(extra_dirs=[settings.fonts_dir])
    cases = build_cases(library)
    pptx = args.out_dir / "calibration.pptx"
    build_pptx(cases, pptx)
    print(f"built {pptx} ({len(cases)} cases)")
    subprocess.run(
        [sys.executable, str(Path(__file__).parent / "render_pptx.py"), str(pptx),
         "--out-dir", str(args.out_dir / "render")],
        check=True,
    )
    measure(cases, args.out_dir / "render" / "render_001.png")


if __name__ == "__main__":
    main()
