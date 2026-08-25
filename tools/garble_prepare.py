#!/usr/bin/env python3
"""注記なしで、変換時のreview(要確認行)から文字化け修正タスクを生成する。

対象: 変換で生成された編集可能PPTX + その work ディレクトリ。
やること(すべて決定論):
  1. work/pages/*/layout.json のreviewから「テキストとして修正可能な行」を抽出
     (低信頼OCR行、行端の未認識インク残存行。画像のまま保持の行は
      テキストが無いため対象外として一覧のみ報告)
  2. 各行が属するブロック(段落群)を特定し、出力PPTX内の該当図形を解決
  3. 除去前の原本画像(source.png)から該当箇所の切り抜き(2倍拡大)を生成
  4. feedback_apply.py と同じスキーマの tasks.json を書き出す

エージェントの仕事は answers.json を書くことだけ。適用は
tools/feedback_apply.py を使う(ハッシュ証跡・本文不変性の検証つき)。
"""
from __future__ import annotations

import argparse
import glob
import hashlib
import json
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

import cv2  # noqa: E402
from pptx import Presentation  # noqa: E402


def iter_shapes_deep(shapes):
    for s in shapes:
        yield s
        if s.shape_type == 6:
            yield from iter_shapes_deep(s.shapes)


def paragraphs_text(shape) -> list[str]:
    return ["".join(r.text for r in p.runs) for p in shape.text_frame.paragraphs]


FIXABLE_REASON_KEYS = ("未認識のインクが残存",)
IMAGE_KEPT_KEYS = ("画像のまま保持",)


def register_tasks(out: "Path", digest: str, payload: dict, kind: str) -> None:
    """自動取り込み(answers_autoingest)用にタスクセットを台帳登録する。"""
    reg = Path(__file__).resolve().parent.parent / "runs" / "tasks_registry.jsonl"
    reg.parent.mkdir(parents=True, exist_ok=True)
    import json as _json
    with open(reg, "a", encoding="utf-8") as f:
        f.write(_json.dumps({
            "sha": digest, "tasks": str((out / "tasks.json").resolve()),
            "source": payload["source"], "work_dir": payload.get("work_dir"),
            "kind": kind}, ensure_ascii=False) + "\n")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=Path, help="編集可能化済みPPTX")
    ap.add_argument("--work-dir", type=Path, default=ROOT / "work")
    ap.add_argument("--out-dir", type=Path, default=ROOT / "feedback_work")
    args = ap.parse_args()

    prs = Presentation(str(args.pptx))
    out = args.out_dir
    crops = out / "crops"
    crops.mkdir(parents=True, exist_ok=True)

    tasks = []
    image_kept = []
    for lp in sorted(glob.glob(str(args.work_dir / "pages" / "*" / "layout.json"))):
        lay = json.loads(Path(lp).read_text("utf-8"))
        page = lay["number"]
        if page > len(prs.slides):
            continue
        cw, ch = lay["canvas_size"]
        flagged = []
        for r in lay.get("review", []):
            reason = r.get("reason", "")
            text = r.get("text", "")
            if not text.strip() or r.get("resolved"):
                continue
            if any(k in reason for k in IMAGE_KEPT_KEYS):
                image_kept.append((page, text[:24], reason))
                continue
            if reason and not any(k in reason for k in FIXABLE_REASON_KEYS):
                continue  # 情報系(修正済み記録など)は対象外
            flagged.append(r)
        if not flagged:
            continue
        # ブロック特定: flagged行のテキストを含むブロックを探す
        slide = prs.slides[page - 1]
        shape_paras = {}
        for s in iter_shapes_deep(slide.shapes):
            if getattr(s, "has_text_frame", False) and s.text_frame.text.strip():
                shape_paras[s.shape_id] = paragraphs_text(s)
        src_img = cv2.imread(str(Path(lp).parent / "source.png"))
        seen_shapes = set()
        for n, r in enumerate(flagged):
            text = r["text"]
            hit = None
            for sid, paras in shape_paras.items():
                if any(text == p for p in paras):
                    hit = sid
                    break
            if hit is None:
                # 完全一致が無い場合は部分一致で試す
                for sid, paras in shape_paras.items():
                    if any(text[:10] and text[:10] in p for p in paras):
                        hit = sid
                        break
            entry = {
                "id": f"g{page:03d}_{n}",
                "slide": page,
                "class": "A",
                "bbox": [round(v, 1) for v in r["bbox"]] if r.get("bbox") else None,
                "marker_text": r.get("reason") or "低信頼OCR行",
                "marker_shape_id": None,
                "target": None,
                "suggested": None,
                "status": "open",
            }
            if hit is not None:
                if hit in seen_shapes:
                    continue  # 同一図形は1タスクに統合
                seen_shapes.add(hit)
                entry["target"] = {"shape_id": hit, "name": "",
                                   "paragraphs": shape_paras[hit]}
            else:
                entry["status"] = "needs_human"
                entry["note"] = "出力PPTX内で該当図形を特定できず"
            # 原本画像からの切り抜き(判断の最重要根拠)
            if src_img is not None and r.get("bbox"):
                x0, y0, x1, y1 = r["bbox"]
                pad = max(8, int((y1 - y0) * 0.6))
                sx0 = max(0, int(x0) - pad)
                sy0 = max(0, int(y0) - pad)
                sx1 = min(cw, int(x1) + pad)
                sy1 = min(ch, int(y1) + pad)
                if sx1 - sx0 > 8 and sy1 - sy0 > 8:
                    crop = cv2.resize(src_img[sy0:sy1, sx0:sx1], (0, 0),
                                      fx=2, fy=2, interpolation=cv2.INTER_LANCZOS4)
                    cp = crops / f"{entry['id']}_source.png"
                    cv2.imwrite(str(cp), crop)
                    entry["source_crop"] = str(cp.relative_to(out))
            tasks.append(entry)

    payload = {
        "source": str(args.pptx),
        "source_sha256": hashlib.sha256(args.pptx.read_bytes()).hexdigest(),
        "work_dir": str(args.work_dir) if args.work_dir else None,
        "tasks": tasks,
    }
    body = json.dumps(payload, ensure_ascii=False, indent=1)
    (out / "tasks.json").write_text(body, "utf-8")
    digest = hashlib.sha256(body.encode()).hexdigest()
    (out / "tasks.sha256").write_text(digest, "utf-8")
    register_tasks(out, digest, payload, "garble")
    print(f"tasks.json: {len(tasks)}件 -> {out / 'tasks.json'}")
    if image_kept:
        print(f"参考: 画像のまま保持されている行 {len(image_kept)}件"
              "(テキストが無いためこのフローでは修正不可):")
        for pg, tx, rs in image_kept[:8]:
            print(f"  p{pg} {tx!r}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
