#!/usr/bin/env python3
"""answers.json を検証してPPTXへ決定論的に適用する。

エージェントの自由度は answers.json の中身だけ。本スクリプトは:
  1. answers.json のスキーマと tasks.json のハッシュ証跡を検証(fail-loud)
  2. status=fix のタスクだけ、指し先図形の「テキストのみ」を差し替える
     (段落数一致を要求。スタイルは各段落の先頭ランを継承)
  3. 注記図形は既定でそのまま残す(--resolve-markers で本文に ✔ を付記)
  4. 出力を再読込検証し、適用結果を runs/ にJSONLで記録する

使い方:
  feedback_apply.py 入力.pptx feedback_work/tasks.json answers.json -o 出力.pptx

answers.json 形式:
  {"tasks_sha256": "<tasks.jsonのsha256>",
   "answers": [{"id": "s007_m21", "status": "fix|skip|needs_human",
                "paragraphs": ["修正後の段落1", "…"]}]}
"""
from __future__ import annotations

import argparse
import datetime
import hashlib
import json
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation  # noqa: E402


def fail(msg: str) -> int:
    print(f"適用拒否: {msg}")
    return 1


def set_paragraph_text(paragraph, new_text: str) -> None:
    """段落のテキストを差し替える。

    文字数が同じ場合は既存ランの区切りで配り直す(色ラン分割等の
    ラン単位の書式を保持)。文字数が変わる場合は先頭ランの書式を継承。
    """
    runs = paragraph.runs
    if not runs:
        paragraph.text = new_text
        return
    old = "".join(r.text for r in runs)
    if len(old) == len(new_text):
        pos = 0
        for r in runs:
            n = len(r.text)
            r.text = new_text[pos:pos + n]
            pos += n
        return
    runs[0].text = new_text
    for r in runs[1:]:
        r.text = ""


def iter_shapes_deep(shapes):
    for s in shapes:
        yield s
        if s.shape_type == 6:  # GROUP
            yield from iter_shapes_deep(s.shapes)


def writeback_workdir(work_dir: Path, slide_no: int,
                      before: list[str], after: list[str]) -> str:
    """適用済み修正を変換キャッシュへ書き戻す。

    layout.jsonの該当ブロック(行テキストが完全一致)を更新し、見つからなければ
    native_*.xml内の段落テキストを更新する。あわせて、修正で変わった行に
    一致するreviewエントリへ resolved を立てる(タスクの再提示防止)。
    戻り値は書き戻し先の種別("layout" / "native" / "none")。
    """
    page_dir = work_dir / "pages" / f"{slide_no:03d}"
    lay_path = page_dir / "layout.json"
    # 修正台帳へ永続記録: フル再解析(バージョンアップ等)でlayout.jsonが
    # 作り直されても、変換パイプラインが合成前にこの台帳をリプレイする
    ledger = work_dir / "fixes_ledger.jsonl"
    entry = json.dumps({"slide": slide_no, "before": before, "after": after},
                       ensure_ascii=False)
    existing = ledger.read_text("utf-8").splitlines() if ledger.is_file() else []
    if entry not in existing:
        with open(ledger, "a", encoding="utf-8") as f:
            f.write(entry + "\n")
    kind = "none"
    if lay_path.is_file():
        lay = json.loads(lay_path.read_text("utf-8"))
        for block in lay.get("blocks", []):
            texts = [ln["text"] for ln in block["lines"]]
            if texts == before:
                for ln, new in zip(block["lines"], after):
                    ln["text"] = new
                kind = "layout"
                break
        changed_lines = {old for old, new in zip(before, after) if old != new}
        for r in lay.get("review", []):
            if r.get("text") in changed_lines:
                r["resolved"] = True
        if kind == "layout" or changed_lines:
            lay_path.write_text(json.dumps(lay, ensure_ascii=False, indent=1),
                                "utf-8")
    if kind == "none":
        # ネイティブ図形が指し先だった場合: native_*.xml の段落を更新
        from lxml import etree
        A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        for xml_path in sorted(page_dir.glob("native_*.xml")):
            root = etree.fromstring(xml_path.read_bytes())
            dirty = False
            paras = [p for p in root.iter(A + "p")]
            for p in paras:
                ts = list(p.iter(A + "t"))
                joined = "".join(t.text or "" for t in ts)
                for old, new in zip(before, after):
                    if old != new and joined == old:
                        if len(new) == len(old):
                            pos = 0
                            for t in ts:
                                n = len(t.text or "")
                                t.text = new[pos:pos + n]
                                pos += n
                        elif ts:
                            ts[0].text = new
                            for t in ts[1:]:
                                t.text = ""
                        dirty = True
            if dirty:
                xml_path.write_text(etree.tostring(root, encoding="unicode"),
                                    "utf-8")
                kind = "native"
                break
    return kind


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("pptx", type=Path)
    ap.add_argument("tasks", type=Path)
    ap.add_argument("answers", type=Path)
    ap.add_argument("-o", "--output", type=Path, required=True)
    ap.add_argument("--resolve-markers", action="store_true",
                    help="適用済みタスクの注記本文の先頭に ✔ を付ける")
    ap.add_argument("--work-dir", type=Path, default=None,
                    help="変換時のworkディレクトリ。指定すると修正をキャッシュ"
                         "(layout.json / native_*.xml)へも書き戻し、reviewを"
                         "解決済みにする。再変換しても修正が保たれ、"
                         "同じタスクが再提示されなくなる")
    args = ap.parse_args()

    tasks_body = args.tasks.read_text("utf-8")
    tasks_sha = hashlib.sha256(tasks_body.encode()).hexdigest()
    tasks = {t["id"]: t for t in json.loads(tasks_body)["tasks"]}
    payload = json.loads(args.answers.read_text("utf-8"))

    # --- ハッシュ証跡: prepare を経ていない/改変された tasks では進めない ---
    if payload.get("tasks_sha256") != tasks_sha:
        return fail("answers.json の tasks_sha256 が tasks.json と一致しません"
                    " (prepare→回答→apply の順序が守られていません)")
    src_sha = hashlib.sha256(args.pptx.read_bytes()).hexdigest()
    if json.loads(tasks_body).get("source_sha256") != src_sha:
        return fail("入力PPTXが tasks.json 生成時と異なります")

    answers = payload.get("answers")
    if not isinstance(answers, list):
        return fail("answers が配列ではありません")

    prs = Presentation(str(args.pptx))
    log = []
    applied = skipped = held = rejected = 0
    # 同一図形に複数タスクを順に適用できるよう、適用後の期待本文を追跡する
    expected: dict[tuple[int, int], list[str]] = {}
    for a in answers:
        aid = a.get("id")
        status = a.get("status")
        t = tasks.get(aid)
        if t is None:
            rejected += 1
            log.append({"id": aid, "result": "rejected", "reason": "未知のタスクID"})
            continue
        if status == "skip":
            skipped += 1
            log.append({"id": aid, "result": "skip"})
            continue
        if status == "needs_human" or t["class"] == "C":
            held += 1
            log.append({"id": aid, "result": "needs_human"})
            continue
        if status != "fix":
            rejected += 1
            log.append({"id": aid, "result": "rejected", "reason": f"不正status: {status}"})
            continue
        new_paras = a.get("paragraphs")
        target = t.get("target")
        if not target or not isinstance(new_paras, list):
            rejected += 1
            log.append({"id": aid, "result": "rejected", "reason": "targetまたはparagraphs欠落"})
            continue
        slide = prs.slides[t["slide"] - 1]
        shape = next((s for s in iter_shapes_deep(slide.shapes)
                      if s.shape_id == target["shape_id"]), None)
        if shape is None or not getattr(shape, "has_text_frame", False):
            rejected += 1
            log.append({"id": aid, "result": "rejected", "reason": "指し先図形が見つからない"})
            continue
        paras = shape.text_frame.paragraphs
        cur = ["".join(r.text for r in p.runs) for p in paras]
        key = (t["slide"], target["shape_id"])
        baseline = expected.get(key, target["paragraphs"])
        if cur != baseline:
            rejected += 1
            log.append({"id": aid, "result": "rejected",
                        "reason": "指し先の本文が tasks 生成時から変わっている"})
            continue
        if len(new_paras) != len(paras):
            rejected += 1
            log.append({"id": aid, "result": "rejected",
                        "reason": f"段落数不一致 ({len(paras)}→{len(new_paras)})"})
            continue
        for p, txt in zip(paras, new_paras):
            if not isinstance(txt, str):
                rejected += 1
                log.append({"id": aid, "result": "rejected", "reason": "段落が文字列でない"})
                break
        else:
            for p, txt in zip(paras, new_paras):
                set_paragraph_text(p, txt)
            expected[key] = list(new_paras)
            applied += 1
            wb = "none"
            if args.work_dir is not None:
                wb = writeback_workdir(args.work_dir, t["slide"], cur, new_paras)
            log.append({"id": aid, "result": "applied",
                        "before": cur, "after": new_paras, "writeback": wb})
            if args.resolve_markers:
                mk = next((s for s in iter_shapes_deep(slide.shapes)
                           if s.shape_id == t["marker_shape_id"]), None)
                if mk is not None and getattr(mk, "has_text_frame", False):
                    p0 = mk.text_frame.paragraphs[0]
                    if p0.runs and not p0.runs[0].text.startswith("✔"):
                        p0.runs[0].text = "✔ " + p0.runs[0].text

    prs.save(str(args.output))
    Presentation(str(args.output))  # 再読込検証(壊れたzipなら例外)

    # --- 学習: 適用された修正から再利用可能な知識を自動抽出する ---
    # 1文字の漢字置換 → 形近字ペア(以後どの資料でも安全ゲートつき候補になる)
    # それ以外の変更 → フレーズ辞書(完全一致置換。同テンプレの資料で即修正)
    learned_pairs: set[tuple[str, str]] = set()
    learned_phrases: dict[str, str] = {}
    kanji = __import__("re").compile(r"[一-鿿々]")
    for row in log:
        if row.get("result") != "applied":
            continue
        for old, new in zip(row.get("before", []), row.get("after", [])):
            if old == new or not old.strip():
                continue
            if len(old) == len(new):
                diffs = [(a, b) for a, b in zip(old, new) if a != b]
                if all(kanji.match(a) and kanji.match(b) for a, b in diffs):
                    learned_pairs.update(diffs)
                    continue
            if len(old) >= 4:
                learned_phrases[old] = new
    rules_dir = ROOT / "rules"
    if learned_pairs or learned_phrases:
        cpath = rules_dir / "learned_confusables.json"
        existing_pairs = set()
        if cpath.is_file():
            existing_pairs = {tuple(p) for p in json.loads(
                cpath.read_text("utf-8")).get("pairs", [])}
        existing_pairs |= learned_pairs
        cpath.write_text(json.dumps(
            {"pairs": sorted([list(p) for p in existing_pairs])},
            ensure_ascii=False, indent=1), "utf-8")
        ppath = rules_dir / "learned_phrases.json"
        phrases = {}
        if ppath.is_file():
            phrases = json.loads(ppath.read_text("utf-8"))
        phrases.update(learned_phrases)
        ppath.write_text(json.dumps(phrases, ensure_ascii=False, indent=1),
                         "utf-8")
        print(f"学習: 形近字ペア+{len(learned_pairs)} フレーズ+{len(learned_phrases)}"
              f" (累計 ペア{len(existing_pairs)} / フレーズ{len(phrases)})")

    runs_dir = ROOT / "runs"
    runs_dir.mkdir(exist_ok=True)
    stamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    log_path = runs_dir / f"feedback_{stamp}.jsonl"
    with open(log_path, "w", encoding="utf-8") as f:
        for row in log:
            f.write(json.dumps(row, ensure_ascii=False) + "\n")
    print(f"適用 {applied} / skip {skipped} / 要人間 {held} / 拒否 {rejected}")
    print(f"出力: {args.output}\nログ: {log_path}")
    return 0 if rejected == 0 else 2


if __name__ == "__main__":
    raise SystemExit(main())
