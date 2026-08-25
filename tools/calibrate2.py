#!/usr/bin/env python3
"""較正プローブv2: 各レンダラの行レイアウト特性をPDF座標から精密測定する。

v1(インク位置のラスタ測定)と違い、エクスポートPDFのテキストスパン原点
(=ベースライン座標)を直接読むため、測定誤差がほぼゼロ。
全ケース2行構成なので「1行目ベースライン位置」と「行送り」を同時に測れる。
ケースIDを行テキスト自体に埋め込むため、どのレンダラのPDFでも照合できる。
フォント置換もPDF埋め込みフォント名から自動検出する。

使い方:
  build:   .venv/bin/python tools/calibrate2.py build
           → work/calib2/probe2.pptx を生成(決定論的。metaファイル不要)
  measure: .venv/bin/python tools/calibrate2.py measure <エクスポート済みPDF>
           → ケースごとの実測 行送り/固定値 比・1行目ベースラインを表で出力
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from notebooklm_pdf2pptx.fontlib import FontLibrary
from notebooklm_pdf2pptx.pptx_writer import (
    SLIDE_H_EMU, SLIDE_W_EMU, _set_exact_line_spacing, _set_run_fonts,
    _disable_kerning, first_baseline_offset_pt, natural_line_height_pt,
)
from notebooklm_pdf2pptx.config import Settings

CANVAS_W, CANVAS_H = 1376, 774  # 960x540ptと同アスペクト(必須)
PT_PER_PX = 960.0 / CANVAS_W
JP = "行送り検証テキスト"
EN = "BASELINE CHECK TEXT"


def build_cases() -> list[dict]:
    """決定論的なケース表。build と measure の両方がこれを使う。

    mode: "exact"=spcPts固定行送り(value=pt) / "pct"=spcPct(value=%)
          / "none"=lnSpc指定なし(レンダラ既定の1行送りを実測する)
    """
    cases = []

    def add(family, bold, size, mode, value):
        cid = f"C{len(cases) + 1:02d}"
        text = JP if "Noto" in family or "Hiragino" in family else EN
        cases.append(dict(id=cid, family=family, bold=bold, size=size,
                          mode=mode, value=value, text=text))

    # 行送り倍率スイープ (Noto Regular 20pt, spcPts = size×r)
    for r in (1.0, 1.2, 1.35, 1.5, 1.7, 2.0):
        add("Noto Sans JP", False, 20.0, "exact", round(20.0 * r, 2))
    # サイズスイープ (r=1.35固定)
    for s in (12.0, 18.0, 26.0, 34.0):
        add("Noto Sans JP", False, s, "exact", round(s * 1.35, 2))
    # フォント差 (メトリクスの大きく違うOswald / 太字)
    add("Oswald", False, 20.0, "exact", 27.0)
    add("Oswald", True, 20.0, "exact", 34.0)
    add("Noto Sans JP", True, 28.0, "exact", 37.8)
    # レンダラ既定の1行送り(lnSpcなし)
    add("Noto Sans JP", False, 20.0, "none", 0.0)
    add("Oswald", False, 20.0, "none", 0.0)
    # パーセント指定
    add("Noto Sans JP", False, 20.0, "pct", 150.0)

    # 配置: 2列×4行グリッドを必要なだけのスライドに展開
    col_x = (40.0, 700.0)
    row_h = 176.0
    for i, c in enumerate(cases):
        slot = i % 8
        c["slide"] = i // 8
        c["left_px"] = col_x[slot % 2]
        c["top_px"] = 30.0 + (slot // 2) * row_h
    return cases


def build(out_dir: Path) -> Path:
    settings = Settings()
    library = FontLibrary(extra_dirs=[settings.fonts_dir])
    presentation = Presentation()
    presentation.slide_width = Emu(SLIDE_W_EMU)
    presentation.slide_height = Emu(SLIDE_H_EMU)
    emu_x = SLIDE_W_EMU / CANVAS_W
    emu_y = SLIDE_H_EMU / CANVAS_H
    cases = build_cases()
    slides = []
    for _ in range(max(c["slide"] for c in cases) + 1):
        slides.append(presentation.slides.add_slide(presentation.slide_layouts[6]))

    for c in cases:
        face = library.face(c["family"], bold=c["bold"])
        if face is None:
            raise SystemExit(f"フォントが見つかりません: {c['family']}")
        box = slides[c["slide"]].shapes.add_textbox(
            Emu(int(c["left_px"] * emu_x)), Emu(int(c["top_px"] * emu_y)),
            Emu(int(620 * emu_x)), Emu(int(150 * emu_y)))
        frame = box.text_frame
        frame.clear()
        frame.margin_left = frame.margin_right = 0
        frame.margin_top = frame.margin_bottom = 0
        frame.word_wrap = False
        frame.auto_size = MSO_AUTO_SIZE.NONE
        frame.vertical_anchor = MSO_ANCHOR.TOP
        for li, tag in enumerate(("a", "b")):
            paragraph = frame.paragraphs[0] if li == 0 else frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_before = paragraph.space_after = Pt(0)
            if c["mode"] == "exact":
                _set_exact_line_spacing(paragraph, c["value"])
            elif c["mode"] == "pct":
                paragraph.line_spacing = c["value"] / 100.0
            run = paragraph.add_run()
            run.text = f"{c['id']}{tag} {c['text']}"
            run.font.size = Pt(c["size"])
            run.font.bold = face.bind_bold
            run.font.color.rgb = RGBColor(0, 0, 0)
            _set_run_fonts(run, face.typeface)
            _disable_kerning(run)
    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / "probe2.pptx"
    presentation.save(str(path))
    print(f"built {path} ({len(cases)} cases, {len(slides)} slides)")
    return path


def measure(pdf_path: Path) -> None:
    import fitz

    doc = fitz.open(str(pdf_path))
    meta = doc.metadata or {}
    print(f"PDF creator={meta.get('creator')!r} producer={meta.get('producer')!r}")
    # ページサイズをプローブ座標系(960pt幅)へ正規化
    spans: dict[str, dict] = {}
    for page in doc:
        scale = 960.0 / page.rect.width
        for blk in page.get_text("dict")["blocks"]:
            if blk.get("type") != 0:
                continue
            for line in blk["lines"]:
                for sp in line["spans"]:
                    t = sp["text"].strip()
                    if len(t) >= 4 and t[0] == "C" and t[1:3].isdigit() and t[3] in "ab":
                        spans[t[:4]] = dict(
                            x=sp["origin"][0] * scale, y=sp["origin"][1] * scale,
                            font=sp["font"], size=sp["size"] * scale,
                            page=page.number)

    settings = Settings()
    library = FontLibrary(extra_dirs=[settings.fonts_dir])
    print(f"{'case':30s} {'PDFフォント':22s} {'行送り実測':>9s} {'/指定':>6s} "
          f"{'先頭bl実測':>9s} {'差(pp)':>7s} {'差(lo)':>7s}")
    for c in build_cases():
        a = spans.get(c["id"] + "a")
        b = spans.get(c["id"] + "b")
        label = (f"{c['id']} {c['family'][:10]}{'B' if c['bold'] else ''} "
                 f"{c['size']:g}pt {c['mode']}"
                 + (f"={c['value']:g}" if c["mode"] != "none" else ""))
        if a is None or b is None:
            print(f"{label:30s} スパン未検出")
            continue
        face = library.face(c["family"], bold=c["bold"])
        expect_sub = face.typeface.replace(" ", "")
        sub = "" if expect_sub.lower()[:4] in a["font"].replace("-", "").lower().lower() \
            else " ←置換!"
        pitch = b["y"] - a["y"]
        ratio = pitch / c["value"] if c["mode"] == "exact" else float("nan")
        fbo_meas = a["y"] - (c["top_px"] + c["slide"] * 0) * PT_PER_PX
        if c["mode"] == "exact":
            pitch_for_model = c["value"]
        else:
            pitch_for_model = natural_line_height_pt(face.path, face.index, c["size"])
            if c["mode"] == "pct":
                pitch_for_model *= c["value"] / 100.0
        pred = {p: first_baseline_offset_pt(face.path, face.index, c["size"],
                                            pitch_for_model, profile=p)
                for p in ("powerpoint", "libreoffice")}
        print(f"{label:30s} {a['font'][:20] + sub:22s} {pitch:9.2f} "
              f"{ratio:6.3f} {fbo_meas:9.2f} "
              f"{fbo_meas - pred['powerpoint']:+7.2f} "
              f"{fbo_meas - pred['libreoffice']:+7.2f}")
        dx = a["x"] - c["left_px"] * PT_PER_PX
        if abs(dx) > 1.5:
            print(f"{'':30s} 注意: 左端が指定と{dx:+.1f}ptずれ(インセット差?)")


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    sub = parser.add_subparsers(dest="cmd", required=True)
    b = sub.add_parser("build")
    b.add_argument("--out-dir", type=Path, default=Path("work/calib2"))
    m = sub.add_parser("measure")
    m.add_argument("pdf", type=Path)
    args = parser.parse_args()
    if args.cmd == "build":
        build(args.out_dir)
    else:
        measure(args.pdf)


if __name__ == "__main__":
    main()
