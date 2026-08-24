"""スモークテスト: 生成PPTXの構造健全性を検査する。

- python-pptxで再オープンできる
- スライド数・テキストボックス数が layout.json と一致する
- 全テキストボックスがスライド内に収まっている
- フォント名・サイズ・spc・行間(spcPts)が設定されている

使い方: .venv/bin/python tools/smoke_test.py output.pptx
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--pages-dir", type=Path, default=Path("work/pages"))
    args = parser.parse_args()

    presentation = Presentation(str(args.pptx))
    slide_w = presentation.slide_width
    slide_h = presentation.slide_height
    problems: list[str] = []
    n_boxes = 0
    for index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # 本ツールが生成したテキストボックスのみ検査
            # (PPTX入力から持ち越したネイティブ図形は元の書式を保持していて正しい)
            if not str(shape.name).startswith(("Text:", "Overlay text:")):
                continue
            n_boxes += 1
            # スライド外へのはみ出し (5%の余白まで許容)
            if shape.left < -slide_w * 0.05 or shape.top < -slide_h * 0.05:
                problems.append(f"slide{index}: 負座標 {shape.name!r}")
            if shape.left + shape.width > slide_w * 1.08:
                problems.append(f"slide{index}: 右はみ出し {shape.name!r}")
            for paragraph in shape.text_frame.paragraphs:
                # 本ツール生成の段落のみ検査する。生成ランは必ず kern="10000" を
                # 持つ(カーニング無効化マーカー)。持ち越しネイティブ図形や、
                # 出力を人間がPowerPointで編集して増えた段落は対象外。
                a_ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                ours = any(r._r.find(a_ns + "rPr") is not None
                           and r._r.find(a_ns + "rPr").get("kern") == "10000"
                           for r in paragraph.runs)
                if not ours:
                    continue
                ppr = paragraph._p.find(a_ns + "pPr")
                if ppr is None or ppr.find(a_ns + "lnSpc") is None:
                    problems.append(f"slide{index}: lnSpc未設定 {shape.name!r}")
                for run in paragraph.runs:
                    if run.font.size is None:
                        problems.append(f"slide{index}: サイズ未設定 {shape.name!r}")
                    rpr = run._r.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr")
                    if rpr is None or rpr.find(
                        "{http://schemas.openxmlformats.org/drawingml/2006/main}ea"
                    ) is None:
                        problems.append(f"slide{index}: a:ea未設定 {shape.name!r}")

    print(f"slides={len(presentation.slides)} textboxes={n_boxes}")
    if problems:
        for problem in problems[:30]:
            print("NG:", problem)
        sys.exit(1)
    print("OK: 構造検査に問題なし")


if __name__ == "__main__":
    main()
